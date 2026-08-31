import json
import os
from pathlib import Path
from typing import Dict, List, Tuple

from config import get_logger
from data_layer.datalayer_exceptions.datalayer_exceptions import InvalidFileType

logger = get_logger(__name__)

BINARY_SNIFF_BYTES = 8192
BYTE_ORDER_MARKS = (
    (b"\xef\xbb\xbf", "utf-8-sig"),
    (b"\xff\xfe\x00\x00", "utf-32"),
    (b"\x00\x00\xfe\xff", "utf-32"),
    (b"\xff\xfe", "utf-16"),
    (b"\xfe\xff", "utf-16"),
)

OFFICE_BINARY_EXTENSIONS = {".doc", ".odt", ".odp", ".ods", ".rtf", ".epub", ".xls"}


class TextExtractor:
    """Turns a file of any type into plain text.

    Extensions with a structure worth exploiting (PDF, Office, HTML, JSON,
    notebooks) get a dedicated reader; everything else — source code, logs,
    configuration, TeX, subtitles, formats nobody has thought of yet — is
    decoded as text. Only bytes that are not text at all are rejected.
    """

    def __init__(self):
        self._handlers = {
            ".pdf": self._extract_from_pdf,
            ".docx": self._extract_from_docx,
            ".pptx": self._extract_from_pptx,
            ".xlsx": self._extract_from_xlsx,
            ".xlsm": self._extract_from_xlsx,
            ".html": self._extract_from_markup,
            ".htm": self._extract_from_markup,
            ".xhtml": self._extract_from_markup,
            ".xml": self._extract_from_markup,
            ".json": self._extract_from_json,
            ".ipynb": self._extract_from_notebook,
        }
        for extension in OFFICE_BINARY_EXTENSIONS:
            self._handlers[extension] = self._extract_from_binary_document

    def _read_bytes(self, file_path: str) -> bytes:
        with open(file_path, "rb") as file:
            return file.read()

    def _decode(self, raw: bytes, file_path: str) -> str:
        """Decode bytes to text, or raise if they are not text at all.

        chardet is consulted only after the common encodings fail, because it
        guesses from statistics and will happily label a binary blob as some
        obscure codepage.
        """
        # A byte order mark is the only trustworthy signal for the wide
        # encodings: utf-16 accepts almost any even-length input, so trying it
        # blind turns a Latin-1 document into CJK without raising anything.
        for mark, encoding in BYTE_ORDER_MARKS:
            if raw.startswith(mark):
                try:
                    return raw.decode(encoding)
                except UnicodeDecodeError:
                    break

        if b"\x00" in raw[:BINARY_SNIFF_BYTES]:
            raise InvalidFileType(Path(file_path).suffix or "binary content")

        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            pass

        try:
            import chardet

            detected = chardet.detect(raw[:BINARY_SNIFF_BYTES])
            if detected.get("encoding") and detected.get("confidence", 0) >= 0.7:
                return raw.decode(detected["encoding"], errors="replace")
        except (ImportError, LookupError, UnicodeDecodeError):
            logger.debug(f"Could not detect the encoding of '{file_path}'.")

        return raw.decode("latin-1", errors="replace")

    def _extract_from_txt(self, file_path: str) -> str:
        return self._decode(self._read_bytes(file_path), file_path)

    def _extract_from_docx(self, file_path: str) -> str:
        try:
            from docx import Document

            doc = Document(file_path)
            blocks = []
            for paragraph in doc.paragraphs:
                # Word's heading styles are the document's structure; rendering
                # them as markdown is what lets the normalizer see sections in a
                # format that has no visible heading syntax of its own.
                level = self._docx_heading_level(paragraph)
                blocks.append(
                    f"\n{'#' * level} {paragraph.text}\n"
                    if level and paragraph.text.strip()
                    else paragraph.text
                )
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        blocks.append(" | ".join(cells))
            return "\n".join(blocks)
        except ImportError:
            logger.warning("python-docx is not installed. Run: pip install python-docx")
            return ""
        except Exception as e:
            logger.error(f"Error reading DOCX file '{file_path}': {e}", exc_info=True)
            return ""

    def _docx_heading_level(self, paragraph) -> int:
        style = getattr(paragraph.style, "name", "") or ""
        if style == "Title":
            return 1
        if style.startswith("Heading"):
            _, _, tail = style.partition(" ")
            return int(tail) + 1 if tail.isdigit() else 2
        return 0

    def _extract_from_binary_document(self, file_path: str) -> str:
        """Legacy Office, OpenDocument, RTF and EPUB, via textract."""
        try:
            import textract

            return textract.process(file_path).decode("utf-8", errors="replace")
        except ImportError:
            logger.warning(
                f"textract is not installed, so '{file_path}' cannot be read. "
                "Run: pip install textract"
            )
            return ""
        except Exception as e:
            logger.error(f"Error reading '{file_path}': {e}", exc_info=True)
            return ""

    def _extract_from_pdf(self, file_path: str) -> str:
        try:
            import PyPDF2

            text: List[str] = []
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    # extract_text() returns None for pages that are pure image
                    # or have no text layer, which "\n".join then choked on.
                    text.append(page.extract_text() or "")
            return "\n".join(text)
        except ImportError:
            logger.warning("PyPDF2 is not installed. Run: pip install PyPDF2")
            return ""
        except Exception as e:
            logger.error(f"Error reading PDF file '{file_path}': {e}", exc_info=True)
            return ""

    def _extract_from_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation

            blocks: List[str] = []
            for slide in Presentation(file_path).slides:
                # shapes.title builds a fresh proxy on each access, so the
                # title is matched by id rather than by identity.
                title = slide.shapes.title
                title_id = title.shape_id if title is not None else None
                for shape in slide.shapes:
                    if not shape.has_text_frame or not shape.text_frame.text.strip():
                        continue
                    text = shape.text_frame.text
                    blocks.append(
                        f"# {text}" if shape.shape_id == title_id else text
                    )
            return "\n\n".join(blocks)
        except ImportError:
            logger.warning("python-pptx is not installed. Run: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"Error reading PPTX file '{file_path}': {e}", exc_info=True)
            return ""

    def _extract_from_xlsx(self, file_path: str) -> str:
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            blocks: List[str] = []
            for sheet in workbook.worksheets:
                blocks.append(sheet.title.upper())
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) for cell in row if cell is not None]
                    if cells:
                        blocks.append(" | ".join(cells))
            workbook.close()
            return "\n".join(blocks)
        except ImportError:
            logger.warning("openpyxl is not installed. Run: pip install openpyxl")
            return self._extract_from_binary_document(file_path)
        except Exception as e:
            logger.error(f"Error reading spreadsheet '{file_path}': {e}", exc_info=True)
            return ""

    def _extract_from_markup(self, file_path: str) -> str:
        """HTML and XML with the tags removed.

        Markup used to be read as plain text, so every tag and inline script was
        chunked and embedded alongside the prose it surrounds.
        """
        markup = self._extract_from_txt(file_path)
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(markup, "html.parser")
            for element in soup(["script", "style", "noscript"]):
                element.decompose()
            for level in range(1, 7):
                for element in soup.find_all(f"h{level}"):
                    element.replace_with(
                        f"\n\n{'#' * level} {element.get_text(strip=True)}\n\n"
                    )
            return soup.get_text(separator="\n").strip()
        except ImportError:
            logger.warning(
                "beautifulsoup4 is not installed, so markup is kept as raw text. "
                "Run: pip install beautifulsoup4"
            )
            return markup
        except Exception as e:
            logger.error(f"Error parsing markup '{file_path}': {e}", exc_info=True)
            return markup

    def _flatten_json(self, value, prefix: str = "") -> List[str]:
        """One "a.b.c: value" line per leaf, so the keys stay attached to what
        they label once the document is split into chunks."""
        if isinstance(value, dict):
            lines: List[str] = []
            for key, item in value.items():
                lines.extend(
                    self._flatten_json(item, f"{prefix}.{key}" if prefix else str(key))
                )
            return lines
        if isinstance(value, list):
            lines = []
            for index, item in enumerate(value):
                lines.extend(self._flatten_json(item, f"{prefix}[{index}]"))
            return lines
        return [f"{prefix}: {value}".strip() if prefix else str(value)]

    def _extract_from_json(self, file_path: str) -> str:
        raw = self._extract_from_txt(file_path)
        try:
            return "\n".join(line for line in self._flatten_json(json.loads(raw)) if line)
        except (json.JSONDecodeError, RecursionError) as e:
            logger.warning(f"'{file_path}' is not valid JSON ({e}); read as plain text.")
            return raw

    def _extract_from_notebook(self, file_path: str) -> str:
        raw = self._extract_from_txt(file_path)
        try:
            notebook = json.loads(raw)
        except json.JSONDecodeError as e:
            logger.warning(f"'{file_path}' is not a readable notebook ({e}).")
            return raw

        blocks: List[str] = []
        for cell in notebook.get("cells", []):
            source = cell.get("source", [])
            text = "".join(source) if isinstance(source, list) else str(source)
            if text.strip():
                blocks.append(text)
        return "\n\n".join(blocks)

    def extract_text_from_file(self, file_path) -> Tuple[str, str]:
        path = str(file_path)
        if not os.path.exists(path):
            logger.error(f"extract_text_from_file failed: File does not exist: '{path}'")
            raise FileNotFoundError(path)

        extension = Path(path).suffix.lower()
        logger.debug(f"Extracting text from '{path}' (extension: {extension})")

        handler = self._handlers.get(extension, self._extract_from_txt)
        return path, handler(path)

    def extract_all(self, loaded_files: Dict[str, list]) -> Dict[str, str]:
        extracted_texts: Dict[str, str] = {}
        total_files = sum(len(paths) for paths in loaded_files.values())
        logger.info(
            f"Starting batch text extraction across {total_files} file(s) in {len(loaded_files)} categories..."
        )

        for category, file_paths in loaded_files.items():
            if file_paths:
                logger.info(f"Extracting category: '{category}' ({len(file_paths)} files)")
            for file_path in file_paths:
                logger.debug(f"Processing extraction for: '{file_path}'")
                try:
                    path, text = self.extract_text_from_file(file_path)
                except (InvalidFileType, FileNotFoundError, OSError) as e:
                    # One unreadable file should not abandon the rest of the batch.
                    logger.warning(f"Skipped '{file_path}': {e}")
                    continue
                if text.strip():
                    extracted_texts[path] = text
                else:
                    logger.warning(f"Skipped '{file_path}': no text could be extracted.")

        logger.info(f"Successfully extracted text from {len(extracted_texts)} file(s).")
        return extracted_texts
